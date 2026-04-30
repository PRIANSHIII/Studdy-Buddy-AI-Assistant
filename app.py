import os
import json
import base64
import datetime
from functools import wraps
from collections import defaultdict
from flask import (
    Flask, render_template, redirect, url_for,
    session, request, jsonify, send_from_directory, flash
)
from dotenv import load_dotenv
from werkzeug.utils import secure_filename
from google import genai
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# --- Load environment variables ---
load_dotenv()

app = Flask(__name__)
app.secret_key = os.getenv('SECRET_KEY', 'your_very_secret_fallback_key')

# ==========================
# SECTION 0: GEMINI CLIENT SETUP & PROGRESS TRACKING
# ==========================
GEMINI_CLIENT = None
api_key = os.getenv("GEMINI_API_KEY")

if not api_key:
    print("\n--- WARNING: GEMINI_API_KEY NOT FOUND ---")
    print("AI features (notes generation, flashcards, chat, TTS) will not work.\n")
else:
    try:
        GEMINI_CLIENT = genai.Client(api_key=api_key)
        print("Gemini client initialized successfully.")
    except Exception as e:
        print(f"Error initializing Gemini Client: {e}")

# Progress tracking
USER_PROGRESS = defaultdict(lambda: {
    'study_sessions': [],
    'topics_covered': [],
    'flashcards_studied': 0,
    'quizzes_taken': 0,
    'last_active': None
})


# ==========================
# SECTION 1: AUTHENTICATION
# ==========================

def requires_auth(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get('user_authenticated'):
            return redirect(url_for('login'))
        return f(*args, **kwargs)

    return decorated


@app.route('/')
def home():
    if session.get('user_authenticated'):
        return redirect(url_for('study_buddy'))
    return redirect(url_for('login'))


@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        session['user_authenticated'] = True
        session['user_name'] = 'TestUser'
        session['user_email'] = 'test@example.com'
        return jsonify({'success': True, 'redirect': url_for('study_buddy')})
    return render_template('login.html')


@app.route('/logout')
def logout():
    session.pop('user_authenticated', None)
    session.pop('user_name', None)
    session.pop('user_email', None)
    return redirect(url_for('login'))


# ==========================
# SECTION 2: MAIN PAGES
# ==========================

@app.route('/study-buddy')
@requires_auth
def study_buddy():
    return render_template('index.html', user_name=session.get('user_name', 'Student'))


@app.route('/profile', methods=['GET', 'POST'])
@requires_auth
def profile():
    profile_data = {
        'name': session.get('user_name', 'Student User'),
        'age': 'N/A',
        'email': session.get('user_email', 'test@example.com'),
        'class': 'Not Set',
        'institution': 'Not Set'
    }
    if request.method == 'POST':
        profile_data.update(request.form)
        session['user_name'] = profile_data['name']
        return redirect(url_for('profile'))
    return render_template('profile.html', profile=profile_data)


# ==========================
# SECTION 3: SYLLABUS UPLOAD
# ==========================

UPLOAD_FOLDER = 'uploads'
SYLLABUS_FOLDER = os.path.join(UPLOAD_FOLDER, 'syllabus')
NOTES_UPLOAD_FOLDER = os.path.join(UPLOAD_FOLDER, 'notes')
GENERATED_NOTES_FOLDER = os.path.join('generated_notes')

# Ensure folders exist
os.makedirs(SYLLABUS_FOLDER, exist_ok=True)
os.makedirs(NOTES_UPLOAD_FOLDER, exist_ok=True)
os.makedirs(GENERATED_NOTES_FOLDER, exist_ok=True)

ALLOWED_EXTENSIONS = {'pdf', 'doc', 'docx', 'txt', 'ppt', 'pptx'}


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def allowed_note_file(filename):
    return allowed_file(filename)


@app.route('/syllabus', methods=['GET', 'POST'])
@requires_auth
def syllabus():
    uploaded_files = os.listdir(SYLLABUS_FOLDER)
    uploaded_files = [f for f in uploaded_files if not f.startswith('.')]

    if request.method == 'POST':
        if 'syllabus_file' not in request.files:
            return redirect(request.url)
        file = request.files['syllabus_file']
        if file.filename == '':
            return redirect(request.url)
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file.save(os.path.join(SYLLABUS_FOLDER, filename))
            return redirect(url_for('syllabus'))
        else:
            return render_template('syllabus.html', files=uploaded_files,
                                   error="Invalid file type. Only PDF, DOCX, DOC, or TXT are allowed.")
    return render_template('syllabus.html', files=uploaded_files)


@app.route('/uploads/syllabus/<filename>')
@requires_auth
def uploaded_file(filename):
    return send_from_directory(SYLLABUS_FOLDER, filename)


@app.route('/delete_file', methods=['POST'])
@requires_auth
def delete_file():
    filename = request.form.get('filename')
    if filename:
        file_path_syllabus = os.path.join(SYLLABUS_FOLDER, filename)
        file_path_notes = os.path.join(NOTES_UPLOAD_FOLDER, filename)
        file_path_generated = os.path.join(GENERATED_NOTES_FOLDER, filename)
        for path in [file_path_syllabus, file_path_notes, file_path_generated]:
            if os.path.exists(path):
                os.remove(path)
                flash(f"{filename} deleted successfully.", "success")
    return redirect(request.referrer or url_for('syllabus'))


@app.route('/generate_study_plan', methods=['POST'])
@requires_auth
def generate_study_plan():
    filename = request.form.get('filename')
    file_path = os.path.join(SYLLABUS_FOLDER, filename)
    if not os.path.exists(file_path):
        flash("File not found!", "error")
        return redirect(url_for('syllabus'))

    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            syllabus_text = f.read()
    except Exception as e:
        flash(f"Failed to read file: {e}", "error")
        return redirect(url_for('syllabus'))

    if not GEMINI_CLIENT:
        flash("AI core isn't running. Check your setup.", "error")
        return redirect(url_for('syllabus'))

    try:
        system_instruction = (
            "You are 'Study Buddy', an AI assistant that generates structured study plans "
            "from syllabus content. Keep it concise, organized, and practical for exam preparation."
        )
        prompt = f"Create a detailed study plan for this syllabus:\n\n{syllabus_text}"

        response = GEMINI_CLIENT.models.generate_content(
            model='gemini-2.5-flash',
            contents=[{"text": prompt}],
            config={"system_instruction": system_instruction}
        )

        study_plan_text = response.text
        plan_filename = f"study_plan_for_{filename.rsplit('.', 1)[0]}.txt"
        plan_path = os.path.join(SYLLABUS_FOLDER, plan_filename)
        with open(plan_path, 'w', encoding='utf-8') as f:
            f.write(study_plan_text)

        flash(f"Study plan generated for {filename}! Download as {plan_filename}", "success")
        return redirect(url_for('syllabus'))
    except Exception as e:
        flash(f"Failed to generate study plan: {e}", "error")
        return redirect(url_for('syllabus'))


# ==========================
# SECTION 4: NOTES UPLOAD & GENERATION
# ==========================

def extract_text(file_path):
    ext = file_path.rsplit('.', 1)[1].lower()
    text = ""

    if ext == 'pdf':
        reader = PdfReader(file_path)
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
    elif ext in ['doc', 'docx']:
        doc = Document(file_path)
        text = "\n".join(p.text for p in doc.paragraphs)
    elif ext in ['ppt', 'pptx']:
        ppt = Presentation(file_path)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif ext == 'txt':
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
    return text


@app.route('/notes', methods=['GET', 'POST'])
@requires_auth
def notes():
    # Ensure folders exist
    os.makedirs(NOTES_UPLOAD_FOLDER, exist_ok=True)
    os.makedirs(GENERATED_NOTES_FOLDER, exist_ok=True)

    uploaded_notes = [f for f in os.listdir(NOTES_UPLOAD_FOLDER) if not f.startswith('.')]
    generated_notes = [f for f in os.listdir(GENERATED_NOTES_FOLDER) if not f.startswith('.')]

    if request.method == 'POST' and 'notes_file' in request.files:
        file = request.files['notes_file']
        if file.filename != '':
            filename = secure_filename(file.filename)
            file.save(os.path.join(NOTES_UPLOAD_FOLDER, filename))
            flash(f"{filename} uploaded successfully!", "success")
            return redirect(url_for('notes'))

    return render_template('notes.html', uploaded_notes=uploaded_notes, generated_notes=generated_notes)


@app.route('/generate_notes', methods=['POST'])
@requires_auth
def generate_notes():
    filename = request.form.get('filename')
    file_path = os.path.join(NOTES_UPLOAD_FOLDER, filename)

    if not os.path.exists(file_path):
        flash("File not found!", "error")
        return redirect(url_for('notes'))

    if not GEMINI_CLIENT:
        flash("AI core isn't running. Check your setup.", "error")
        return redirect(url_for('notes'))

    try:
        # Correct function call
        text_content = extract_text(file_path)

        system_instruction = (
            "You are 'Study Buddy', an AI assistant that generates well-structured, clear, "
            "modern, and concise study notes from uploaded notes content."
        )

        response = GEMINI_CLIENT.models.generate_content(
            model='gemini-2.5-flash',
            contents=[{"text": text_content}],
            config={"system_instruction": system_instruction}
        )

        generated_text = response.text
        gen_filename = f"{filename.rsplit('.', 1)[0]}_notes.txt"
        gen_path = os.path.join(GENERATED_NOTES_FOLDER, gen_filename)

        with open(gen_path, 'w', encoding='utf-8') as f:
            f.write(generated_text)

        flash(f"Generated notes for {filename}!", "success")
        return redirect(url_for('notes'))

    except Exception as e:
        flash(f"Failed to generate notes: {e}", "error")
        return redirect(url_for('notes'))


@app.route('/download/generated_notes/<filename>')
@requires_auth
def download_generated_notes(filename):
    return send_from_directory(GENERATED_NOTES_FOLDER, filename, as_attachment=True)


@app.route('/delete_notes_file', methods=['POST'])
@requires_auth
def delete_notes_file():
    filename = request.form.get('filename')
    path = os.path.join(NOTES_UPLOAD_FOLDER, filename)
    if os.path.exists(path):
        os.remove(path)
        flash(f"Uploaded note {filename} deleted successfully.", "success")
    else:
        flash("File not found.", "error")
    return redirect(url_for('notes'))


@app.route('/delete_generated_note', methods=['POST'])
@requires_auth
def delete_generated_note():
    filename = request.form.get('filename')
    path = os.path.join(GENERATED_NOTES_FOLDER, filename)
    if os.path.exists(path):
        os.remove(path)
        flash(f"Generated note {filename} deleted successfully.", "success")
    else:
        flash("File not found.", "error")
    return redirect(url_for('notes'))


@app.route('/generate_flashcards_ajax', methods=['POST'])
@requires_auth
def generate_flashcards_ajax():
    data = request.get_json()
    filename = data.get('filename')
    file_path = os.path.join(GENERATED_NOTES_FOLDER, filename)

    if not os.path.exists(file_path):
        return jsonify({"success": False, "message": "Generated note file not found."})

    if not GEMINI_CLIENT:
        return jsonify({"success": False, "message": "AI core isn't running. Check your setup."})

    try:
        # 1️⃣ Read the generated note
        with open(file_path, 'r', encoding='utf-8') as f:
            note_text = f.read()

        # 2️⃣ Instruction for AI - improved prompt
        system_instruction = (
            "You are 'Study Buddy', an AI assistant that converts study notes into "
            "concise, clear, structured flashcards. "
            "Create 5-10 flashcards from the provided notes. "
            "Return ONLY valid JSON format: [{'question': 'question text', 'answer': 'answer text'}] "
            "Do not include any other text, explanations, or formatting outside the JSON array."
        )

        prompt = f"Convert these study notes into flashcards:\n\n{note_text}\n\nReturn only JSON array."

        # 3️⃣ Call Gemini API
        response = GEMINI_CLIENT.models.generate_content(
            model='gemini-2.5-flash',
            contents=[{"text": prompt}],
            config={"system_instruction": system_instruction}
        )

        raw_text = response.text.strip()
        print("Raw AI Response:", raw_text)

        # 4️⃣ Clean the response and parse JSON
        # Remove any markdown code blocks if present
        cleaned_text = raw_text
        if '```json' in cleaned_text:
            cleaned_text = cleaned_text.split('```json')[1].split('```')[0]
        elif '```' in cleaned_text:
            cleaned_text = cleaned_text.split('```')[1].split('```')[0]

        # Try to parse JSON
        try:
            flashcards = json.loads(cleaned_text)
            if not isinstance(flashcards, list):
                raise ValueError("JSON is not a list")

            # Validate each flashcard has required fields
            validated_flashcards = []
            for card in flashcards:
                if isinstance(card, dict) and 'question' in card and 'answer' in card:
                    validated_flashcards.append({
                        'question': str(card['question']).strip(),
                        'answer': str(card['answer']).strip()
                    })

            if not validated_flashcards:
                raise ValueError("No valid flashcards found in response")

            print(f"Successfully generated {len(validated_flashcards)} flashcards")
            return jsonify({"success": True, "flashcards": validated_flashcards})

        except json.JSONDecodeError as e:
            print(f"JSON parsing failed: {e}")
            # Fallback: try to extract Q&A pairs from text
            flashcards = extract_flashcards_from_text(raw_text)
            if flashcards:
                return jsonify({"success": True, "flashcards": flashcards})
            else:
                return jsonify({"success": False, "message": "Failed to parse flashcards from AI response"})

    except Exception as e:
        print(f"Flashcards generation error: {e}")
        return jsonify({"success": False, "message": f"Error generating flashcards: {str(e)}"})


def extract_flashcards_from_text(text):
    """Fallback method to extract flashcards from plain text"""
    flashcards = []
    lines = text.split('\n')

    current_question = None
    current_answer = None

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # Look for Q: pattern
        if line.lower().startswith('q:') or line.startswith('Question:'):
            # Save previous card if exists
            if current_question and current_answer:
                flashcards.append({
                    'question': current_question,
                    'answer': current_answer
                })
            current_question = line.split(':', 1)[1].strip() if ':' in line else line[2:].strip()
            current_answer = None

        # Look for A: pattern
        elif line.lower().startswith('a:') or line.startswith('Answer:'):
            if current_question:
                current_answer = line.split(':', 1)[1].strip() if ':' in line else line[2:].strip()

    # Don't forget the last card
    if current_question and current_answer:
        flashcards.append({
            'question': current_question,
            'answer': current_answer
        })

    return flashcards


@app.route('/uploads/notes/<filename>')
@requires_auth
def uploaded_notes_file(filename):
    if os.path.exists(os.path.join(NOTES_UPLOAD_FOLDER, filename)):
        return send_from_directory(NOTES_UPLOAD_FOLDER, filename)
    elif os.path.exists(os.path.join(GENERATED_NOTES_FOLDER, filename)):
        return send_from_directory(GENERATED_NOTES_FOLDER, filename)
    else:
        flash("File not found", "error")
        return redirect(url_for('notes'))


# ==========================
# SECTION 5: AI CHAT & TTS
# ==========================

@app.route('/api/chat', methods=['POST'])
@requires_auth
def chat_with_buddy():
    data = request.json
    user_message = data.get('message')
    if not user_message:
        return jsonify({'error': 'No message provided'}), 400

    if not GEMINI_CLIENT:
        return jsonify({'response': "AI core isn't running. Check your setup."})

    try:
        system_instruction = (
            "You are 'Study Buddy', a friendly AI assistant helping students. "
            "Keep responses concise, encouraging, and highly relevant to study topics."
        )
        config = {"system_instruction": system_instruction}

        response = GEMINI_CLIENT.models.generate_content(
            model='gemini-2.5-flash',
            contents=user_message,  # Fixed: pass the string directly
            config=config
        )
        ai_response = response.text

        # Track progress
        track_progress_internal({
            'activity': 'chat_interaction',
            'duration': 0,
            'topic': 'AI Chat'
        })

        return jsonify({'response': ai_response})
    except Exception as e:
        print(f"Gemini API Error: {e}")
        return jsonify({'response': "Error processing request. Please try again."}), 500


@app.route('/api/tts', methods=['POST'])
@requires_auth
def generate_tts():
    data = request.json
    text = data.get('text')
    if not text:
        return jsonify({'error': 'No text provided for TTS'}), 400
    if not GEMINI_CLIENT:
        return jsonify({'error': "AI core offline."}), 503

    try:
        tts_config = {
            "responseModalities": ["AUDIO"],
            "speechConfig": {
                "voiceConfig": {"prebuiltVoiceConfig": {"voiceName": "Kore"}}
            }
        }

        response = GEMINI_CLIENT.models.generate_content(
            model='gemini-2.5-flash-preview-tts',
            contents=[{"parts": [{"text": text}]}],
            config=tts_config
        )

        if (response.candidates and
                response.candidates[0].content and
                response.candidates[0].content.parts and
                hasattr(response.candidates[0].content.parts[0], 'inline_data') and
                response.candidates[0].content.parts[0].inline_data):

            inline_data = response.candidates[0].content.parts[0].inline_data
            audio_bytes = inline_data.data
            mime_type = inline_data.mime_type
            audio_base64 = base64.b64encode(audio_bytes).decode('utf-8')
            return jsonify({'audioData': audio_base64, 'mimeType': mime_type})
        else:
            return jsonify({'error': 'TTS generation failed. Audio missing.'}), 500
    except Exception as e:
        print(f"TTS API Error: {e}")
        return jsonify({'error': f"TTS API error: {e}"}), 500


# ==========================
# SECTION 6: PROGRESS TRACKING
# ==========================

def track_progress_internal(data):
    user_id = session.get('user_email')
    USER_PROGRESS[user_id]['study_sessions'].append({
        'timestamp': datetime.datetime.now().isoformat(),
        'activity': data.get('activity'),
        'duration': data.get('duration', 0),
        'topic': data.get('topic')
    })
    USER_PROGRESS[user_id]['last_active'] = datetime.datetime.now().isoformat()


@app.route('/api/track_progress', methods=['POST'])
@requires_auth
def track_progress():
    data = request.json
    user_id = session.get('user_email')

    USER_PROGRESS[user_id]['study_sessions'].append({
        'timestamp': datetime.datetime.now().isoformat(),
        'activity': data.get('activity', 'unknown'),
        'duration': data.get('duration', 0)
    })

    if data.get('topic'):
        USER_PROGRESS[user_id]['topics_covered'].append(data['topic'])

    USER_PROGRESS[user_id]['last_active'] = datetime.datetime.now().isoformat()

    return jsonify({'success': True})


@app.route('/api/start_study_session', methods=['POST'])
@requires_auth
def start_study_session():
    session['study_session_start'] = datetime.datetime.now().isoformat()
    session['current_topic'] = request.json.get('topic')
    return jsonify({'success': True})


@app.route('/api/end_study_session', methods=['POST'])
@requires_auth
def end_study_session():
    if session.get('study_session_start'):
        start_time = datetime.datetime.fromisoformat(session['study_session_start'])
        duration = (datetime.datetime.now() - start_time).total_seconds() / 60  # in minutes

        # Track progress
        track_progress_internal({
            'activity': 'study_session',
            'duration': duration,
            'topic': session.get('current_topic')
        })

        session.pop('study_session_start', None)
        session.pop('current_topic', None)

    return jsonify({'success': True})


def calculate_weekly_data(sessions):
    """Calculate study time for each day of the current week"""
    today = datetime.datetime.now()
    start_of_week = today - datetime.timedelta(days=today.weekday())

    weekly_data = [0] * 7  # Initialize with zeros for each day

    for session in sessions:
        session_date = datetime.datetime.fromisoformat(session['timestamp'])
        if session_date >= start_of_week:
            day_index = session_date.weekday()
            weekly_data[day_index] += session.get('duration', 0)

    return weekly_data


@app.route('/progress')
@requires_auth
def progress():
    user_id = session.get('user_email')
    user_progress = USER_PROGRESS.get(user_id, {
        'study_sessions': [],
        'topics_covered': [],
        'flashcards_studied': 0,
        'quizzes_taken': 0,
        'current_streak': 0,
        'last_active': None
    })

    # Calculate metrics
    total_study_time = sum(session.get('duration', 0) for session in user_progress.get('study_sessions', []))
    topics_covered = len(set(user_progress.get('topics_covered', [])))

    # Calculate daily average
    if user_progress.get('study_sessions'):
        days_active = len(set(session['timestamp'][:10] for session in user_progress['study_sessions']))
        daily_average = total_study_time / max(days_active, 1)
    else:
        daily_average = 0

    # Activity counts for chart
    activity_counts = {
        'study_sessions': len(
            [s for s in user_progress.get('study_sessions', []) if s.get('activity') == 'study_session']),
        'chat_interactions': len(
            [s for s in user_progress.get('study_sessions', []) if s.get('activity') == 'chat_interaction']),
        'flashcards': user_progress.get('flashcards_studied', 0),
        'quizzes': user_progress.get('quizzes_taken', 0)
    }

    # Weekly data for chart
    weekly_data = calculate_weekly_data(user_progress.get('study_sessions', []))
    weekly_labels = ['Mon', 'Tue', 'Wed', 'Thu', 'Fri', 'Sat', 'Sun']

    return render_template('progress.html',
                           progress=user_progress,
                           total_study_time=total_study_time,
                           topics_covered=topics_covered,
                           daily_average=daily_average,
                           activity_counts=activity_counts,
                           weekly_data=weekly_data,
                           weekly_labels=weekly_labels)


@app.route('/api/health')
@requires_auth
def health_check():
    return jsonify({
        'connected': True,
        'ai_available': GEMINI_CLIENT is not None,
        'timestamp': datetime.datetime.now().isoformat()
    })


# ==========================
# SECTION 7: RUN APP
# ==========================

if __name__ == '__main__':
    if not os.getenv("GEMINI_API_KEY"):
        print("\n--- WARNING: GEMINI_API_KEY NOT FOUND ---")
        print("AI chat and TTS will not function until the key is set.\n")
    app.run(debug=True, port=5006)